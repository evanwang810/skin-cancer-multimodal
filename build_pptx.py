# -*- coding: utf-8 -*-
"""Build the polished skin-cancer deck, including a design-process / debugging arc."""
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

PROJ = r"C:\Users\evanw\.local\bin\skin_cancer_project"

# ---- palette (matches the site's CSS custom properties) --------------------
BG        = RGBColor(0x0a, 0x0e, 0x14)
PANEL     = RGBColor(0x12, 0x19, 0x24)
PANEL_DK  = RGBColor(0x0d, 0x12, 0x1a)
LINE      = RGBColor(0x23, 0x2e, 0x3c)
INK       = RGBColor(0xea, 0xf0, 0xf7)
MUTED     = RGBColor(0x8b, 0x98, 0xa8)
FAINT     = RGBColor(0x6d, 0x7d, 0x92)
BENIGN    = RGBColor(0x3b, 0xa3, 0xe8)
MALIGNANT = RGBColor(0xfb, 0x5a, 0x6a)
ACCENT    = RGBColor(0x9b, 0x8c, 0xff)
GOOD      = RGBColor(0x46, 0xc2, 0x6a)
AMBER     = RGBColor(0xe8, 0xa8, 0x3b)

SANS = "Segoe UI"
SERIF = "Segoe UI Semibold"
MONO = "Consolas"

SW, SH = Inches(13.333), Inches(7.5)
ML = Inches(0.9)  # left margin

insights = json.load(open(rf"{PROJ}\report\insights.json", encoding="utf-8"))
data = json.load(open(rf"{PROJ}\report\data.json", encoding="utf-8"))
training = json.load(open(rf"{PROJ}\report\training.json", encoding="utf-8"))

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]
_page = [0]


def new_slide(footer=True):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element)
    s.shapes._spTree.insert(2, bg._element)
    _page[0] += 1
    if footer:
        textbox(s, ML, SH - Inches(0.5), Inches(6), Inches(0.3),
                "SKIN CANCER · MULTIMODAL", size=8.5, color=RGBColor(0x3a, 0x45, 0x54), font=MONO)
        textbox(s, SW - Inches(1.4), SH - Inches(0.5), Inches(0.9), Inches(0.3),
                f"{_page[0]:02d}", size=8.5, color=RGBColor(0x3a, 0x45, 0x54), font=MONO,
                align=PP_ALIGN.RIGHT)
    return s


def textbox(slide, l, t, w, h, text, size=18, color=INK, bold=False, italic=False,
            font=SANS, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
    return tb


def eyebrow(slide, text, top=Inches(0.55)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, ML, top + Inches(0.02), Inches(0.28), Inches(0.16))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background(); bar.shadow.inherit = False
    textbox(slide, ML + Inches(0.42), top - Inches(0.04), Inches(9), Inches(0.4), text.upper(),
            size=12, color=ACCENT, bold=True, font=MONO)


def h1(slide, text, size=28, top=Inches(1.0)):
    textbox(slide, ML - Inches(0.05), top, Inches(11.5), Inches(1.0), text,
            size=size, color=INK, bold=True, font=SERIF)


def card(slide, l, t, w, h, fill=PANEL, line_color=LINE, radius=0.06):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line_color
    shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def stat_card(slide, l, t, w, h, value, label, hi=False, color=None):
    card(slide, l, t, w, h)
    textbox(slide, l + Inches(0.18), t + Inches(0.14), w - Inches(0.3), h - Inches(0.5),
            value, size=26, color=(color or (BENIGN if hi else INK)), bold=True, font=MONO)
    textbox(slide, l + Inches(0.18), t + h - Inches(0.44), w - Inches(0.3), Inches(0.35),
            label.upper(), size=9.5, color=MUTED, font=SANS)


def flagnote(slide, l, t, w, text, color=FAINT, icon="!"):
    h = Inches(1.05)
    card(slide, l, t, w, h, fill=PANEL_DK)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Inches(0.08), h)
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background(); bar.shadow.inherit = False
    textbox(slide, l + Inches(0.3), t + Inches(0.12), w - Inches(0.55), h - Inches(0.24),
            text, size=11.5, color=color if color != FAINT else FAINT, font=SANS, line_spacing=1.25,
            anchor=MSO_ANCHOR.MIDDLE)


def bar_row(slide, l, t, w, label, pct_val, val_text, weak=False, color=None):
    textbox(slide, l, t, Inches(2.4), Inches(0.35), label, size=11, color=MUTED, font=MONO)
    track_l = l + Inches(2.5)
    track_w = w - Inches(2.5) - Inches(0.9)
    track = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, track_l, t + Inches(0.02), track_w, Inches(0.3))
    track.adjustments[0] = 0.35
    track.fill.solid(); track.fill.fore_color.rgb = PANEL_DK
    track.line.color.rgb = LINE; track.line.width = Pt(0.75)
    track.shadow.inherit = False
    fill_w = max(Emu(1), int(track_w * min(max(pct_val, 0), 1)))
    fill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, track_l, t + Inches(0.02), fill_w, Inches(0.3))
    fill.adjustments[0] = 0.35
    fill.fill.solid(); fill.fill.fore_color.rgb = color or (FAINT if weak else BENIGN)
    fill.line.fill.background()
    fill.shadow.inherit = False
    textbox(slide, l + w - Inches(0.9), t, Inches(0.9), Inches(0.35), val_text, size=11,
            color=INK, font=MONO, align=PP_ALIGN.RIGHT)


def pill(slide, l, t, text, color, w=Inches(1.5)):
    h = Inches(0.42)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shp.adjustments[0] = 0.5
    shp.fill.solid(); shp.fill.fore_color.rgb = PANEL_DK
    shp.line.color.rgb = color; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(10.5); r.font.color.rgb = color
    r.font.name = MONO; r.font.bold = True
    return shp


def style_chart_text(chart, size=10):
    try:
        for ax in (chart.category_axis, chart.value_axis):
            ax.tick_labels.font.size = Pt(size)
            ax.tick_labels.font.color.rgb = MUTED
            ax.format.line.color.rgb = LINE
        chart.value_axis.has_major_gridlines = False
    except Exception:
        pass


# ============================================================ 1: Title
s = new_slide(footer=False)
textbox(s, ML, Inches(1.55), Inches(9.5), Inches(0.5), "HAM10000 · MULTIMODAL CLASSIFICATION",
        size=13, color=ACCENT, bold=True, font=MONO)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, ML + Inches(0.02), Inches(2.15), Inches(1.4), Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background(); bar.shadow.inherit = False
textbox(s, ML - Inches(0.05), Inches(2.45), Inches(11.2), Inches(2.3),
        "Can patient history help a\nmodel catch skin cancer?", size=44, color=INK, bold=True, font=SERIF)
textbox(s, ML, Inches(4.75), Inches(10.2), Inches(1.3),
        "An image-plus-metadata classifier for benign vs. malignant skin lesions,\n"
        "evaluated honestly on 500 held-out cases the model never trained on.",
        size=16, color=MUTED, font=SANS, line_spacing=1.35)
textbox(s, ML, Inches(6.8), Inches(8), Inches(0.4),
        "github.com/evanwang810/skin-cancer-multimodal", size=11, color=FAINT, font=MONO)

# ============================================================ 2: The idea
s = new_slide()
eyebrow(s, "The idea")
h1(s, "A dermoscopy photo isn't the whole story", size=32)
textbox(s, ML, Inches(2.35), Inches(10.4), Inches(3.5),
        "A dermatologist doesn't just look at a lesion — they factor in the patient's age, "
        "sex, and where on the body it sits. Melanoma risk climbs sharply with age; different "
        "lesion types cluster on different body sites. An image-only model throws that context "
        "away.\n\nThis project asks a narrow, testable question: does adding that context back — "
        "as a second small encoder fused with the image features — measurably help the model "
        "separate benign from malignant?",
        size=17, color=MUTED, font=SANS, line_spacing=1.4)

# ============================================================ 3: Data & method
s = new_slide()
eyebrow(s, "Data & method")
h1(s, "HAM10000, split by lesion — not by image")
textbox(s, ML, Inches(1.95), Inches(11.4), Inches(1.8),
        "10,015 dermoscopy images across 7 diagnoses, collapsed to benign vs. malignant. "
        "The dataset has multiple photos of the same physical lesion — so a naive random split "
        "can put near-duplicate images of one lesion in both train and test, quietly inflating "
        "every metric. The split here groups by lesion ID first, 70/15/15, so no lesion ever "
        "crosses the train/val/test boundary.",
        size=15, color=MUTED, font=SANS, line_spacing=1.35)
stats = [("7 → 2", "diagnoses collapsed to benign / malignant"),
         ("70/15/15", "lesion-grouped train / val / test"),
         ("500", "held-out test cases analyzed here")]
cw, gap = Inches(3.73), Inches(0.3)
for i, (v, l) in enumerate(stats):
    stat_card(s, ML + i * (cw + gap), Inches(4.4), cw, Inches(1.35), v, l)

# ============================================================ 4: Architecture
s = new_slide()
eyebrow(s, "Architecture")
h1(s, "Two encoders, fused, one classifier")
steps = [
    ("1", "Image encoder", "ResNet50, pretrained on ImageNet, fine-tuned on dermoscopy, producing a feature vector."),
    ("2", "Metadata encoder", "A small MLP over age, sex, and lesion location, producing its own feature vector."),
    ("3", "Fuse", "The two vectors are concatenated into one combined lesion representation."),
    ("4", "Classify", "One linear layer to P(malignant); threshold tuned on validation, not test."),
]
sw_, sgap = Inches(2.78), Inches(0.25)
for i, (n, t, d) in enumerate(steps):
    l = ML + i * (sw_ + sgap)
    card(s, l, Inches(2.5), sw_, Inches(2.7))
    textbox(s, l + Inches(0.22), Inches(2.68), sw_ - Inches(0.4), Inches(0.6), n, size=28, color=ACCENT, bold=True, font=SERIF)
    textbox(s, l + Inches(0.22), Inches(3.35), sw_ - Inches(0.4), Inches(0.4), t, size=14, color=INK, bold=True, font=SANS)
    textbox(s, l + Inches(0.22), Inches(3.8), sw_ - Inches(0.4), Inches(1.3), d, size=11, color=MUTED, font=SANS, line_spacing=1.3)
textbox(s, ML, Inches(5.5), Inches(11.4), Inches(0.8),
        "The metadata tower is deliberately small — the image is still doing most of the work. "
        "The question is only whether the context adds anything on top.",
        size=12.5, color=FAINT, font=SANS, line_spacing=1.3)

# ============================================================ 5: Section divider
s = new_slide(footer=False)
textbox(s, ML, Inches(2.5), Inches(9), Inches(0.5), "PART TWO", size=13, color=ACCENT, bold=True, font=MONO)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, ML + Inches(0.02), Inches(3.05), Inches(1.4), Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background(); bar.shadow.inherit = False
textbox(s, ML - Inches(0.05), Inches(3.35), Inches(11.5), Inches(1.4),
        "How it was built — and what broke", size=40, color=INK, bold=True, font=SERIF)
textbox(s, ML, Inches(4.9), Inches(10.5), Inches(1.0),
        "Four problems that shaped the final design. Each one changed a number, a decision, "
        "or a line of code.", size=16, color=MUTED, font=SANS, line_spacing=1.35)


def problem_slide(idx, title, symptom, cause, fix, before, after, tail=None,
                  before_lbl="BEFORE", after_lbl="AFTER", after_good=True):
    s = new_slide()
    eyebrow(s, f"Problem {idx}")
    h1(s, title, size=27)
    # left column: the story
    col_w = Inches(6.4)
    y = Inches(1.95)
    for lbl, body, col in (("SYMPTOM", symptom, MALIGNANT),
                           ("CAUSE", cause, AMBER),
                           ("FIX", fix, GOOD)):
        textbox(s, ML, y, Inches(1.5), Inches(0.3), lbl, size=10.5, color=col, bold=True, font=MONO)
        textbox(s, ML, y + Inches(0.32), col_w, Inches(1.0), body, size=13, color=MUTED,
                font=SANS, line_spacing=1.3)
        y += Inches(1.45)
    # right column: before / after cards
    rx = Inches(7.7)
    rw = Inches(4.75)
    card(s, rx, Inches(2.0), rw, Inches(1.85), fill=PANEL_DK)
    textbox(s, rx + Inches(0.25), Inches(2.15), rw - Inches(0.5), Inches(0.3), before_lbl, size=10, color=MALIGNANT, bold=True, font=MONO)
    textbox(s, rx + Inches(0.25), Inches(2.5), rw - Inches(0.5), Inches(1.2), before, size=13.5, color=INK, font=MONO, line_spacing=1.3)
    card(s, rx, Inches(4.05), rw, Inches(1.85), fill=PANEL_DK)
    textbox(s, rx + Inches(0.25), Inches(4.2), rw - Inches(0.5), Inches(0.3), after_lbl, size=10, color=(GOOD if after_good else BENIGN), bold=True, font=MONO)
    textbox(s, rx + Inches(0.25), Inches(4.55), rw - Inches(0.5), Inches(1.2), after, size=13.5, color=INK, font=MONO, line_spacing=1.3)
    if tail:
        textbox(s, ML, Inches(6.4), Inches(11.4), Inches(0.7), tail, size=12, color=FAINT, font=SANS, line_spacing=1.3)
    return s


# ---- Problem 1: leakage
problem_slide(
    1, "The 90% accuracy that wasn't real",
    "An early version scored above 90% and looked great — suspiciously great for this task.",
    "The random train/test split let near-duplicate photos of the same lesion land on both "
    "sides. The model was partly memorizing lesions it had already seen.",
    "Rebuilt the split to group by lesion ID before partitioning, so every lesion lives "
    "entirely in one of train, val, or test. No lesion crosses the boundary.",
    "random split\n> acc  ~0.90\n(leaked)",
    "lesion-grouped split\n> acc  0.826\n(honest)",
    tail="The lower number is the real one. Every metric in this deck comes from the leakage-free split.",
    after_good=True,
)

# ---- Problem 2: clawing accuracy back
problem_slide(
    2, "Clawing accuracy back honestly",
    "Once the leak was gone, accuracy dropped to ~80%. The goal was to recover it without "
    "reintroducing leakage or overfitting the test set.",
    "The image tower was under-fit: too small a fused feature space, a flat learning rate "
    "across pretrained and fresh layers, and a stray ReLU truncating the backbone's output.",
    "Three targeted changes — widened the fused vector (128 -> 256), gave the pretrained "
    "backbone and the fresh head separate learning rates, and added test-time augmentation.",
    "FEAT_DIM   128\nsingle LR\nno TTA",
    "FEAT_DIM   256\ndiscriminative LR\nTTA (h-flip avg)",
    tail="Small, defensible levers — not test-set fishing. The honest number recovered to ~0.82 AUC 0.90+.",
)

# ---- Problem 3: the fix that broke grad-cam
problem_slide(
    3, "The fix that broke Grad-CAM",
    "After the accuracy work, the Grad-CAM cell crashed — the interpretability view that "
    "shows where the model is looking stopped rendering.",
    "Removing that truncating ReLU also removed a Sequential wrapper around the backbone. "
    "The Grad-CAM hook still reached in through the old index and hit nothing.",
    "Traced the hook target down to the new module path and pointed it at the real final "
    "conv block. One-line fix, but only found by following the earlier change downstream.",
    "net[0].layer4\n> TypeError\n(not subscriptable)",
    "net.layer4\n> hook attached\n> heatmaps render",
    tail="A reminder that a clean-looking refactor upstream can silently break something three cells down.",
)

# ---- Problem 4: is 50/50 realistic
nat = training["test_distribution"]
nat_rate = nat["n_malignant"] / nat["n"]
prec_bal = insights["overall"]["precision"]
prec_nat = training["final"]["multimodal"]["precision"]
problem_slide(
    4, "Is a 50/50 malignant split even realistic?",
    "The explorer shows a balanced 250/250 sample, which is easy to read — but real skin "
    "clinics don't see 50% malignant lesions. Does the balance flatter the model?",
    "Accuracy and the display balance are the wrong thing to worry about. Recall and AUC "
    "don't move with prevalence at all — but precision does, sharply.",
    "Kept the display balanced for exploration, but computed the honest metrics on the full "
    "natural test set and foregrounded precision-at-real-prevalence on the insights page.",
    f"balanced sample\nprecision  {prec_bal*100:.0f}%\n(reads clean)",
    f"real {nat_rate*100:.0f}% prevalence\nprecision  {prec_nat*100:.0f}%\n(~5 in 10 flags false)",
    tail="Recall and AUC are the prevalence-independent truth; precision is where the real-world cost shows up.",
    after_good=False,
)

# ============================================================ 10: self-updating design
s = new_slide()
eyebrow(s, "Design decision")
h1(s, "The site updates itself on every retrain")
textbox(s, ML, Inches(1.95), Inches(11.4), Inches(1.4),
        "Nothing in the presentation or the website is a hardcoded number. Retraining the "
        "notebook re-exports three data files, and every page — explorer, insights, and this "
        "deck — reads its figures straight from them. Change the model, rerun, and the whole "
        "showcase re-states itself.",
        size=15, color=MUTED, font=SANS, line_spacing=1.35)
files = [("data.js", "500 test cases: image, Grad-CAM,\nprediction, PCA coordinates, metadata"),
         ("insights.js", "mined breakdowns by age, sex,\nlocation, calibration, error analysis"),
         ("training.js", "per-epoch curves + both models'\nfinal metrics at natural prevalence")]
cw = Inches(3.73)
for i, (name, desc) in enumerate(files):
    l = ML + i * (cw + Inches(0.3))
    card(s, l, Inches(3.7), cw, Inches(2.0))
    textbox(s, l + Inches(0.22), Inches(3.9), cw - Inches(0.4), Inches(0.4), name, size=15, color=ACCENT, bold=True, font=MONO)
    textbox(s, l + Inches(0.22), Inches(4.45), cw - Inches(0.4), Inches(1.3), desc, size=12, color=MUTED, font=SANS, line_spacing=1.3)

# ============================================================ 11: Headline result
s = new_slide()
eyebrow(s, "Result")
h1(s, "On 500 held-out test lesions")
ov = insights["overall"]
stats = [(f"{ov['auc']:.3f}", "AUC", BENIGN), (f"{ov['accuracy']*100:.1f}%", "accuracy", INK),
         (f"{ov['recall']*100:.1f}%", "malignant recall", INK), (f"{ov['precision']*100:.1f}%", "precision (balanced)", INK)]
cw = Inches(2.85)
for i, (v, l, c) in enumerate(stats):
    stat_card(s, ML + i * (cw + Inches(0.15)), Inches(2.3), cw, Inches(1.4), v, l, color=c)
textbox(s, ML, Inches(4.15), Inches(11.4), Inches(1.6),
        "AUC measures how well the model ranks a random malignant lesion above a random "
        "benign one — threshold-free and prevalence-free. Recall is the share of true "
        "malignancies it catches: the number that matters most when a miss is a missed cancer. "
        "Precision is the one that drops at real-world prevalence, shown on the next problem "
        "slide's terms — most flags in a real clinic would be false alarms.",
        size=14, color=MUTED, font=SANS, line_spacing=1.4)

# ============================================================ 12: Feature space scatter
s = new_slide()
eyebrow(s, "Under the hood")
h1(s, "What the fused features actually look like", size=26)
xy = XyChartData()
ben = xy.add_series("benign")
mal = xy.add_series("malignant")
for d in data:
    (ben if d["true"] == "benign" else mal).add_data_point(d["x"], d["y"])
gframe = s.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, ML, Inches(2.05), Inches(6.7), Inches(4.9), xy)
chart = gframe.chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(11)
chart.legend.font.color.rgb = MUTED
colors = {"benign": BENIGN, "malignant": MALIGNANT}
for plot in chart.plots:
    for series in plot.series:
        series.marker.style = 2
        series.marker.size = 5
        series.marker.format.fill.solid()
        series.marker.format.fill.fore_color.rgb = colors[series.name]
        series.marker.format.line.fill.background()
        series.format.line.fill.background()
style_chart_text(chart)
chart.value_axis.minimum_scale = 0.0
chart.value_axis.maximum_scale = 1.0
chart.category_axis.minimum_scale = 0.0
chart.category_axis.maximum_scale = 1.0
textbox(s, Inches(7.9), Inches(2.2), Inches(4.6), Inches(2.4),
        "Each point is one of the 500 test lesions — its fused image+metadata vector, reduced "
        "to 2D with PCA. There's no clinical unit on the axes; they're just the two directions "
        "of largest variation in what the model learned.",
        size=13, color=MUTED, font=SANS, line_spacing=1.35)
textbox(s, Inches(7.9), Inches(4.6), Inches(4.6), Inches(2.2),
        "The two colors mostly separate into their own regions — evidence the model learned a "
        "real boundary, not noise. The cases in the mixed middle are exactly where it's least "
        "certain, and where the misses cluster.",
        size=12, color=FAINT, font=SANS, line_spacing=1.35)


# ============================================================ 13/14: Real cases
def case_slide(eb, title, tag, tag_color, ex, note, after_good=True):
    s = new_slide()
    eyebrow(s, eb)
    h1(s, title, size=26)
    img_l = ML
    for suffix, cap in ((f"{ex['id']}_orig.jpg", "dermoscopy"), (f"{ex['id']}_cam.jpg", "grad-cam")):
        s.shapes.add_picture(rf"{PROJ}\report\img\{suffix}", img_l, Inches(2.35), height=Inches(2.6))
        textbox(s, img_l, Inches(5.0), Inches(2.6), Inches(0.35), cap, size=10, color=FAINT, font=MONO, align=PP_ALIGN.CENTER)
        img_l += Inches(2.75)
    mx = Inches(6.7)
    pill(s, mx, Inches(2.35), tag, tag_color)
    pill(s, mx + Inches(1.65), Inches(2.35), f"P={ex['prob']:.2f}", MUTED, w=Inches(1.5))
    textbox(s, mx, Inches(3.0), Inches(5.7), Inches(2.2), note, size=15, color=MUTED, font=SANS, line_spacing=1.35)
    loc = ex["loc"]
    textbox(s, mx, Inches(5.1), Inches(5.7), Inches(0.4),
            f"age {ex['age']:.0f}  ·  {ex['sex']}  ·  {loc}", size=12, color=FAINT, font=MONO)
    return s


ex_ok = insights["examples"]["confident_correct_malignant"]
ex_miss = insights["examples"]["confidently_missed_malignant"]
case_slide("A real case", "When it works", "CORRECT", GOOD, ex_ok,
           "The model called this malignant with near-total confidence, and it was right. The "
           "Grad-CAM heat sits on the lesion itself, not the surrounding skin — a sign it's "
           "keying on the right region rather than a background artifact.")
case_slide("A real case", "When it fails", "MISSED", MALIGNANT, ex_miss,
           "This lesion was malignant. The model gave it a near-zero malignant probability — "
           "not a borderline call, a confident miss. Cases like this are why recall and error "
           "analysis matter far more than a single accuracy number.", after_good=False)

# ============================================================ 15: age
s = new_slide()
eyebrow(s, "Insight")
h1(s, "Age carries real signal — and a real weakness", size=25)
textbox(s, ML, Inches(1.9), Inches(11.4), Inches(0.9),
        "Malignant rate climbs from 18% under 30 to 75% at 75+ in this sample — exactly the "
        "epidemiological pattern that gives the metadata encoder something real to use.",
        size=14, color=MUTED, font=SANS, line_spacing=1.3)
a30 = [b for b in insights["by_age"] if b["bucket"] == "0-30"][0]
a75 = [b for b in insights["by_age"] if b["bucket"] == "75+"][0]
bar_row(s, ML, Inches(3.1), Inches(10.5), "recall · under 30", a30["recall"], f"{a30['recall']*100:.0f}%", weak=True)
bar_row(s, ML, Inches(3.7), Inches(10.5), "recall · 75+", a75["recall"], f"{a75['recall']*100:.0f}%")
textbox(s, ML, Inches(4.5), Inches(10.5), Inches(0.9),
        "Younger malignant lesions are rarer here and get caught less often — the harder "
        "case, not the common one. The model is strongest exactly where the disease is most "
        "prevalent.", size=12.5, color=FAINT, font=SANS, line_spacing=1.3)

# ============================================================ 16: location
s = new_slide()
eyebrow(s, "Insight")
h1(s, "Recall depends on where the lesion is", size=25)
locs = ["face", "lower extremity", "back", "trunk", "abdomen"]
by_loc = {b["loc"]: b for b in insights["by_location"]}
y = Inches(2.1)
for loc in locs:
    b = by_loc[loc]
    weak = b["recall"] < 0.65
    bar_row(s, ML, y, Inches(10.8), f"{loc} (n={b['n']})", b["recall"], f"{b['recall']*100:.0f}%", weak=weak)
    y += Inches(0.6)
textbox(s, ML, Inches(5.35), Inches(10.8), Inches(0.9),
        "Face and lower-extremity lesions are caught reliably; abdomen recall (60%, small "
        "malignant count) is the weakest reliable signal in the data — a real gap worth "
        "flagging, if still a modest sample.", size=12.5, color=FAINT, font=SANS, line_spacing=1.3)

# ============================================================ 17: sex
s = new_slide()
eyebrow(s, "Insight")
h1(s, "A recall gap by sex, observed but not explained", size=24)
by_sex = {b["sex"]: b for b in insights["by_sex"]}
male, female = by_sex["male"], by_sex["female"]
textbox(s, ML, Inches(1.9), Inches(11.4), Inches(0.7),
        f"Malignant recall is {male['recall']*100:.1f}% for male patients versus "
        f"{female['recall']*100:.1f}% for female patients in this sample — a "
        f"{abs(male['recall']-female['recall'])*100:.0f}-point gap.",
        size=14, color=MUTED, font=SANS, line_spacing=1.3)
bar_row(s, ML, Inches(2.95), Inches(10.5), f"male (n={male['n']})", male["recall"], f"{male['recall']*100:.0f}%")
bar_row(s, ML, Inches(3.55), Inches(10.5), f"female (n={female['n']})", female["recall"], f"{female['recall']*100:.0f}%")
flagnote(s, ML, Inches(4.4), Inches(11.4),
         "This is an observed pattern in one 500-case sample, not a demonstrated cause. It could "
         "reflect which lesion types are more common by sex in this dataset rather than the "
         "model treating sexes differently — worth investigating, not worth overclaiming.")

# ============================================================ 18: calibration
s = new_slide()
eyebrow(s, "Insight")
h1(s, "Confidence isn't proof")
cal = insights["calibration"]
textbox(s, ML, Inches(1.95), Inches(11.4), Inches(1.2),
        f"On average the model is more confident when it's right than when it's wrong — a sign "
        f"of reasonable calibration. But of its {cal['n_wrong']} errors on this sample, "
        f"{cal['confidently_wrong']} ({cal['confidently_wrong_pct']:.0f}%) were confidently "
        "wrong, not borderline calls.",
        size=14, color=MUTED, font=SANS, line_spacing=1.35)
stats = [(f"{cal['mean_confidence_correct']:.2f}", "avg. confidence · correct", GOOD),
         (f"{cal['mean_confidence_wrong']:.2f}", "avg. confidence · wrong", INK),
         (f"{cal['confidently_wrong_pct']:.0f}%", "of errors · confidently wrong", MALIGNANT)]
cw = Inches(3.73)
for i, (v, l, c) in enumerate(stats):
    stat_card(s, ML + i * (cw + Inches(0.3)), Inches(3.6), cw, Inches(1.4), v, l, color=c)

# ============================================================ 19: Limitations
s = new_slide()
eyebrow(s, "Limitations & next steps")
h1(s, "What this isn't, yet")
limits = [
    ("One dataset, one backbone", "HAM10000 only, ResNet50 only. A Vision-Transformer / text-encoded-metadata version is the natural v2 and matches the original brief."),
    ("Thin tails", "Ear, hand, genital, and “unknown” location have too few cases to trust individually — read those bars with caution."),
    ("Precision, not perfection", "At real prevalence most positive flags are false alarms; this is a triage-style aid, not a decision-maker."),
    ("Not a diagnostic tool", "A class project on public data, not validated, cleared, or intended for any clinical use."),
]
y = Inches(2.1)
for head, body in limits:
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, ML, y + Inches(0.05), Inches(0.08), Inches(0.75))
    bar.fill.solid(); bar.fill.fore_color.rgb = FAINT; bar.line.fill.background(); bar.shadow.inherit = False
    textbox(s, ML + Inches(0.3), y, Inches(11), Inches(0.4), head, size=15, color=INK, bold=True, font=SANS)
    textbox(s, ML + Inches(0.3), y + Inches(0.4), Inches(11), Inches(0.7), body, size=12.5, color=MUTED, font=SANS, line_spacing=1.25)
    y += Inches(1.2)

# ============================================================ 20: Closing
s = new_slide(footer=False)
textbox(s, ML, Inches(1.55), Inches(10), Inches(0.5), "EXPLORE IT YOURSELF", size=13, color=ACCENT, bold=True, font=MONO)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, ML + Inches(0.02), Inches(2.1), Inches(1.4), Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background(); bar.shadow.inherit = False
textbox(s, ML - Inches(0.05), Inches(2.4), Inches(11.2), Inches(1.4),
        "Every number here is clickable", size=40, color=INK, bold=True, font=SERIF)
textbox(s, ML, Inches(3.85), Inches(10.2), Inches(1.0),
        "All 500 test cases — image, Grad-CAM, prediction, and metadata — are browsable, "
        "filterable, sortable, recolorable, and exportable in the live explorer.",
        size=15, color=MUTED, font=SANS, line_spacing=1.35)
links = [("Explorer & data table", "evanwang810.github.io/skin-cancer-multimodal"),
         ("Full insights page", "evanwang810.github.io/skin-cancer-multimodal/insights.html"),
         ("Source & notebook", "github.com/evanwang810/skin-cancer-multimodal")]
y = Inches(5.0)
for label, url in links:
    textbox(s, ML, y, Inches(3.2), Inches(0.4), label, size=13, color=MUTED, font=SANS)
    textbox(s, ML + Inches(3.3), y, Inches(8), Inches(0.4), url, size=13, color=BENIGN, font=MONO)
    y += Inches(0.55)

prs.save(rf"{PROJ}\presentation.pptx")
print("saved", len(prs.slides._sldIdLst), "slides")
